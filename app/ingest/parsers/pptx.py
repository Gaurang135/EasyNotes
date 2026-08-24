from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from app.models import ParsedDoc, TextBlock, Table
from app.errors import CorruptFileError, NoExtractableTextError
from app.ingest.ocr import ocr_image, ocr_enabled

_GROUP = 6    # MSO_SHAPE_TYPE.GROUP
_PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE


def _walk(shapes):
    """Yield ('prose', text), ('table', grid), ('image', blob) from a shape tree,
    recursing groups. Tables and grouped shapes hold most of a real deck's content;
    images (diagrams) carry the rest and need OCR."""
    for sh in shapes:
        try:
            if sh.shape_type == _GROUP:
                yield from _walk(sh.shapes)
            elif getattr(sh, "has_table", False):
                grid = [[c.text.strip() for c in row.cells] for row in sh.table.rows]
                grid = [r for r in grid if any(r)]
                if grid:
                    yield ("table", grid)
            elif sh.shape_type == _PICTURE and ocr_enabled():
                try:
                    yield ("image", sh.image.blob)
                except Exception:
                    pass
            elif sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    yield ("prose", t)
        except Exception:
            continue


class PptxParser:
    file_types = frozenset({"pptx"})

    def parse(self, path: Path) -> ParsedDoc:
        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise CorruptFileError(f"unreadable PPTX: {e}")
        blocks: list[TextBlock] = []
        tables: list[Table] = []
        for i, slide in enumerate(prs.slides, 1):
            title = None
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
            except Exception:
                title = None
            prose: list[str] = []
            tnum = 0
            for kind, payload in _walk(slide.shapes):
                if kind == "prose":
                    prose.append(payload)
                elif kind == "image":
                    text = ocr_image(payload)
                    if text:
                        blocks.append(TextBlock(text=text, kind="prose",
                                                location=f"slide {i} diagram (OCR)", heading=title))
                else:  # table
                    tnum += 1
                    rows = ["\t".join(r) for r in payload]
                    blocks.append(TextBlock(text="\n".join(rows), kind="table",
                                            location=f"slide {i} table {tnum}", heading=title))
                    tables.append(Table(name=title or f"Slide {i} table {tnum}",
                                        columns=payload[0], rows=payload[1:],
                                        location=f"slide {i}"))
            if prose:
                blocks.append(TextBlock(text="\n".join(prose), kind="prose",
                                        location=f"slide {i}", heading=title))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(TextBlock(text=notes, kind="prose",
                                            location=f"slide {i} notes", heading=title))
        if not blocks:
            raise NoExtractableTextError("no extractable text")
        return ParsedDoc(text_blocks=blocks, metadata={}, warnings=[], tables=tables)
