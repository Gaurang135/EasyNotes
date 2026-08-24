"""Dev tool: generate ~70 realistic, data-rich documents across all formats.
Usage: python scripts/gen_corpus.py /tmp/corpus
"""
from __future__ import annotations
import sys, random
from pathlib import Path

OUT = Path(sys.argv[1] if len(sys.argv) > 1 else "/tmp/corpus")
OUT.mkdir(parents=True, exist_ok=True)
random.seed(42)

VENDORS = ["Acme Corp", "Globex Ltd", "Initech", "Umbrella Inc", "Soylent Co",
           "Hooli", "Stark Industries", "Wayne Enterprises", "Wonka Foods", "Tyrell Corp"]
REGIONS = ["North", "South", "East", "West", "Central"]
DEPTS = ["Engineering", "Design", "Sales", "Marketing", "Support", "Finance", "HR", "Operations"]
CITIES = ["London", "Mumbai", "Oslo", "Berlin", "Tokyo", "Austin", "Toronto", "Singapore"]
FIRST = ["Priya", "Sam", "Jo", "Alex", "Nina", "Omar", "Lena", "Raj", "Mia", "Ken", "Ava", "Ravi"]
LAST = ["Nair", "Cole", "Patel", "Kim", "Diaz", "Sharma", "Ng", "Rossi", "Khan", "Meyer"]
PRODUCTS = ["Widget", "Gadget", "Gizmo", "Sprocket", "Cog", "Bolt", "Panel", "Sensor", "Module", "Relay"]


def dt(m, d):
    return f"2026-{m:02d}-{d:02d}"


def person():
    return f"{random.choice(FIRST)} {random.choice(LAST)}"


def money(lo, hi):
    return f"{random.randint(lo, hi):,}.{random.randint(0,99):02d}"


# ── INVOICES: 12 (6 PDF, 6 TXT) ──────────────────────────────────────────────
def invoice_text(i):
    v = VENDORS[i % len(VENDORS)]
    sub, tax = random.randint(500, 9000), 0
    tax = round(sub * 0.18)
    lines = "\n".join(f"  {random.choice(PRODUCTS)} x{random.randint(1,20)}  Rs.{money(100,2000)}"
                      for _ in range(random.randint(3, 6)))
    return (f"INVOICE #INV-{2000+i}\nVendor: {v}\nBill to: Client {i}\n"
            f"Invoice date: {dt((i%12)+1, (i%27)+1)}\nDue date: {dt((i%12)+2, (i%27)+1)}\n\n"
            f"Line items:\n{lines}\n\nSubtotal Rs.{sub},000.00\nTax Rs.{tax},000.00\n"
            f"Total Rs.{sub+tax},000.00\n\nContact: billing@{v.split()[0].lower()}.com | "
            f"+91 9{random.randint(100000000,999999999)}\nhttps://{v.split()[0].lower()}.com/pay")


def gen_invoices():
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    for i in range(12):
        body = invoice_text(i)
        if i % 2 == 0:
            p = OUT / f"invoice_{i:02d}.pdf"
            c = canvas.Canvas(str(p), pagesize=letter)
            y = 750
            for line in body.splitlines():
                c.drawString(60, y, line[:95]); y -= 16
                if y < 60:
                    c.showPage(); y = 750
            c.save()
        else:
            (OUT / f"invoice_{i:02d}.txt").write_text(body)


# ── CSVs: 12 (many rows each) ────────────────────────────────────────────────
def gen_csvs():
    def write_csv(name, header, rows):
        (OUT / name).write_text(header + "\n" + "\n".join(rows) + "\n")

    for r in range(4):  # sales_0..3
        rows = [f"{random.choice(REGIONS)},{random.choice(PRODUCTS)},{random.randint(1,500)},"
                f"{random.randint(1000,90000)},{dt((r%12)+1,(k%27)+1)}" for k in range(random.randint(30,60))]
        write_csv(f"sales_{r}.csv", "region,product,units,revenue,date", rows)

    emp = [f"{person()},{random.choice(DEPTS)},{random.randint(600000,2500000)},"
           f"{dt((k%12)+1,(k%27)+1)},{random.choice(CITIES)}" for k in range(45)]
    write_csv("employees.csv", "name,dept,salary,joined,city", emp)

    inv = [f"SKU-{1000+k},{random.choice(PRODUCTS)},{random.randint(0,900)},"
           f"{random.randint(50,5000)},{random.choice(CITIES)}" for k in range(50)]
    write_csv("inventory.csv", "sku,product,qty,price,warehouse", inv)

    for t in range(3):  # transactions
        rows = [f"TXN-{10000+k},{dt((k%12)+1,(k%27)+1)},{random.randint(100,50000)},"
                f"{random.choice(VENDORS)},{random.choice(['paid','pending','failed','refunded'])}"
                for k in range(random.randint(40, 80))]
        write_csv(f"transactions_{t}.csv", "id,date,amount,vendor,status", rows)

    exp = [f"{random.choice(DEPTS)},{random.choice(['travel','software','office','meals'])},"
           f"{random.randint(50,9000)},{dt((k%12)+1,(k%27)+1)}" for k in range(40)]
    write_csv("expenses.csv", "dept,category,amount,date", exp)

    cust = [f"{person()},{random.choice(CITIES)},{random.choice(['free','pro','enterprise'])},"
            f"{random.randint(0,50000)}" for k in range(35)]
    write_csv("customers.csv", "name,city,plan,mrr", cust)

    tickets = [f"TKT-{500+k},{random.choice(['bug','feature','question'])},"
               f"{random.choice(['open','closed','pending'])},{random.choice(FIRST)}" for k in range(30)]
    write_csv("support_tickets.csv", "id,type,status,assignee", tickets)


# ── XLSX: 8 (multi-sheet, many rows) ─────────────────────────────────────────
def gen_xlsx():
    from openpyxl import Workbook
    for i in range(5):
        wb = Workbook()
        rev = wb.active; rev.title = "Revenue"
        rev.append(["month", "region", "revenue", "target"])
        for k in range(24):
            rev.append([dt((k % 12) + 1, 1), random.choice(REGIONS),
                        random.randint(10000, 200000), random.randint(10000, 200000)])
        costs = wb.create_sheet("Costs")
        costs.append(["dept", "amount", "quarter"])
        for k in range(16):
            costs.append([random.choice(DEPTS), random.randint(5000, 90000), f"Q{(k%4)+1}"])
        wb.save(OUT / f"financials_{i}.xlsx")
    for name in ("budget", "payroll", "kpis"):
        wb = Workbook(); ws = wb.active; ws.title = name.title()
        ws.append(["item", "value", "owner"])
        for k in range(20):
            ws.append([f"{name}-{k}", random.randint(1000, 500000), person()])
        wb.save(OUT / f"{name}.xlsx")


# ── PPTX: 8 decks ────────────────────────────────────────────────────────────
def gen_pptx():
    from pptx import Presentation
    titles = ["Q1 Business Review", "Q2 Business Review", "Product Roadmap 2026",
              "Sales Kickoff", "All-Hands Update", "Marketing Plan", "Hiring Plan", "Board Deck"]
    for i, title in enumerate(titles):
        prs = Presentation()
        for s in range(random.randint(3, 6)):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{title} — Section {s+1}"
            slide.placeholders[1].text = (f"Revenue Rs.{random.randint(1,50)}M in {random.choice(REGIONS)}.\n"
                                          f"Owner: {person()}. Target date {dt((i%12)+1,(s%27)+1)}.\n"
                                          f"Contact ops@{VENDORS[i%len(VENDORS)].split()[0].lower()}.com")
            slide.notes_slide.notes_text_frame.text = f"Speaker notes for {title} section {s+1}."
        prs.save(OUT / f"deck_{i:02d}.pptx")


# ── DOCX: 8 reports (headings + tables) ──────────────────────────────────────
def gen_docx():
    from docx import Document
    topics = ["Annual Report", "Incident Postmortem", "Vendor Agreement", "Research Findings",
              "Onboarding Guide", "Security Review", "Quarterly Summary", "Project Charter"]
    for i, topic in enumerate(topics):
        d = Document()
        d.add_heading(f"{topic} {i}", level=1)
        d.add_paragraph(f"Prepared by {person()} on {dt((i%12)+1,(i%27)+1)}. "
                        f"Vendor: {VENDORS[i%len(VENDORS)]}. Budget Rs.{random.randint(1,20)}0,000. "
                        f"Contact: report@corp.io")
        d.add_heading("Summary", level=2)
        d.add_paragraph("Key results and highlights for the period under review.")
        t = d.add_table(rows=1, cols=3)
        t.rows[0].cells[0].text, t.rows[0].cells[1].text, t.rows[0].cells[2].text = "metric", "value", "owner"
        for k in range(6):
            r = t.add_row().cells
            r[0].text, r[1].text, r[2].text = random.choice(PRODUCTS), str(random.randint(1, 9999)), person()
        d.save(OUT / f"report_{i:02d}.docx")


# ── MD + TXT: fill to ~70 ────────────────────────────────────────────────────
def gen_md_txt():
    for i in range(6):
        (OUT / f"spec_{i}.md").write_text(
            f"# Spec {i}: {random.choice(PRODUCTS)} Service\n\n## Overview\n"
            f"Owner {person()}, due {dt((i%12)+1,(i%27)+1)}.\n\n## API\n- GET /{random.choice(PRODUCTS).lower()}\n"
            f"\n## Notes\nBudget Rs.{random.randint(1,9)}0,000. contact dev@corp.io\n")
    notes = [
        ("meeting_q1.txt", "Standup {d}. Attendees: {a}, {b}. Decision: ship v2 by {d2}. contact pm@corp.io"),
        ("resume_1.txt", "{p} — Senior Engineer\nEmail: {e}\nPhone: +91 90000 11111\nSkills: Python, Go, K8s\nJoined 2018"),
        ("resume_2.txt", "{p} — Product Designer\nEmail: {e}\nExperience 6 years in {c}"),
        ("contract_1.txt", "Agreement with Vendor: {v}\nEffective {d}\nValue Rs.5,00,000\nlegal@{vl}.com"),
        ("receipt_1.txt", "Receipt from {v}\nDate {d}\nAmount Rs.1,299.50\nbilling@{vl}.com"),
        ("article_1.txt", "The future of offline document search. Published {d}. Contact editor@news.io. "
                          "Retrieval improved 12% year over year across benchmarks."),
        ("memo_1.txt", "Memo: budget freeze until {d2}. All Rs. approvals via finance@corp.io."),
        ("policy_1.txt", "Refund policy: customers may request a refund within 30 days. support@corp.io"),
    ]
    for name, tmpl in notes:
        (OUT / name).write_text(tmpl.format(
            d=dt(random.randint(1,12), random.randint(1,27)), d2=dt(random.randint(1,12), random.randint(1,27)),
            a=random.choice(FIRST), b=random.choice(FIRST), p=person(),
            e=f"{random.choice(FIRST).lower()}@dev.io", c=random.choice(CITIES),
            v=(vv := random.choice(VENDORS)), vl=vv.split()[0].lower()))


gen_invoices(); gen_csvs(); gen_xlsx(); gen_pptx(); gen_docx(); gen_md_txt()
files = sorted(p.name for p in OUT.iterdir())
print(f"generated {len(files)} files")
from collections import Counter
print(dict(Counter(f.rsplit('.', 1)[-1] for f in files)))
