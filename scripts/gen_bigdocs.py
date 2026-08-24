"""Dev tool: generate a few BIG, legit multi-page PDFs and multi-slide PPTX decks
with coherent real-world content. Usage: python scripts/gen_bigdocs.py /tmp/bigdocs
"""
from __future__ import annotations
import sys
from pathlib import Path

OUT = Path(sys.argv[1] if len(sys.argv) > 1 else "/tmp/bigdocs")
OUT.mkdir(parents=True, exist_ok=True)

# ── Authored content: {title: [(heading, [paragraphs...]), ...]} ──────────────
DOCS: dict[str, list] = {
    "Whitepaper - Offline Document Intelligence": [
        ("Executive Summary", [
            "Organizations accumulate unstructured documents faster than they can organize them: contracts, invoices, decks, spreadsheets, and notes pile up across drives and inboxes. This whitepaper describes an approach to turning that sprawl into clean, queryable data without depending on a large language model at query time.",
            "The system parses each document into text and tables, extracts typed fields with configurable rules, and indexes the content for both keyword and semantic retrieval. The result is a single store that answers precise, database-style questions and open, natural-language questions alike, at zero marginal cost per query."]),
        ("The Problem With Search-Only Tools", [
            "Full-text search finds where words appear; it cannot tell you the total value of last quarter's invoices, filter rows where an amount exceeds a threshold, or match a question to a passage that uses entirely different words. Knowledge-base products solve the natural-language half but forfeit deterministic, field-level querying.",
            "A durable solution needs both modes behind one interface, with a router that sends structured questions to structured data and free-text questions to semantic retrieval."]),
        ("Architecture", [
            "Ingestion is a pipeline: validate, parse, chunk, embed, index. Each format has a dedicated parser so a new format is an additive change. Chunking respects the embedding model's token budget and prepends document context so retrieval stays grounded.",
            "Storage is a single embedded database. Keyword search uses an inverted full-text index with BM25 ranking; semantic search uses a compact sentence-embedding model whose vectors are stored alongside the text. Hybrid queries fuse the two rankings with reciprocal rank fusion, which needs no tuning."]),
        ("Structured Extraction", [
            "Tabular inputs are preserved as typed columns rather than flattened to text, so numeric and date columns remain filterable and sortable. Free-text documents are mined for key-value facts — amounts, dates, emails, phone numbers, and labelled pairs — using a small, extensible rule set inspired by production parser-worker designs.",
            "Because extraction is configuration-driven, adding a new field type is a one-line change, and the same rules run identically across every format."]),
        ("Retrieval Quality", [
            "Several subtle failure modes separate systems that look like they work from systems that do. Vector distance must be converted to similarity consistently, or rankings silently invert. BM25 scores are negative-is-better and must be normalized once. Snippets should anchor on meaningful terms, not stopwords, so a question like 'who is free' does not return every line containing 'is'.",
            "A small labelled evaluation set with recall and mean-reciprocal-rank metrics is the arbiter for every retrieval decision, from chunk size to whether a query instruction helps."]),
        ("Durability and Cost", [
            "The store is a single file, which makes backup and restore trivial. On ephemeral hosts, periodic consistent snapshots to inexpensive object storage provide durability at negligible cost. Embeddings run on CPU, so there is no accelerator dependency and no per-query fee.",
            "This keeps the total running cost near zero while preserving the option to add a generation stage later for composed answers."]),
        ("Security Considerations", [
            "Uploads are validated before parsing: size caps, content sniffing to reject disguised files, and decompression limits to defuse archive bombs. Parsing failures never crash the service; each document records a human-readable reason and the corpus remains queryable.",
            "Original files are retained so the corpus can be re-indexed without re-upload as extraction improves."]),
        ("Conclusion", [
            "Turning messy documents into structured, queryable data is less about any single model and more about disciplined engineering across parsing, extraction, indexing, and retrieval. An LLM-free core keeps the system fast, private, and cheap, while leaving a clean seam for optional answer synthesis.",
            "The approach scales from a personal corpus to a team knowledge base with the same design."]),
    ],
    "Acme FY2026 Annual Report": [
        ("Letter to Shareholders", [
            "Fiscal 2026 was a year of disciplined growth. Revenue rose across every region as our platform expanded from a single product into a suite. We enter the new year with a stronger balance sheet, a broader customer base, and a clear roadmap.",
            "We remained focused on durable unit economics rather than growth at any cost, and that focus is reflected in improving margins and healthy retention."]),
        ("Business Overview", [
            "Acme provides infrastructure that helps companies organize and query their operational data. Our customers range from early-stage teams to large enterprises, and our pricing scales with usage so value and cost stay aligned.",
            "We operate across five regions and serve customers in finance, retail, healthcare, and technology."]),
        ("Financial Highlights", [
            "Total revenue grew to a record level, driven by net new customers and expansion within existing accounts. Gross margin improved as infrastructure efficiency work compounded. Operating expenses grew more slowly than revenue, improving operating leverage.",
            "Free cash flow turned positive for the full year for the first time, giving us flexibility to invest through cycles."]),
        ("Revenue by Region", [
            "Growth was broad-based. The North and West regions led in absolute terms, while Central posted the fastest percentage growth off a smaller base. International revenue continued to increase as a share of the total.",
            "We expect the mix to continue diversifying as our partner ecosystem matures."]),
        ("Product and Technology", [
            "We shipped major improvements to ingestion throughput, retrieval quality, and the analytics experience. Investments in reliability reduced incident frequency and shortened resolution times.",
            "Our research effort focused on retrieval accuracy and cost efficiency, both of which translate directly into customer value."]),
        ("Customers and Markets", [
            "Retention remained strong and expansion accelerated as customers adopted more of the suite. Our net revenue retention reflects both low churn and healthy upsell.",
            "The addressable market continues to expand as more workloads move from manual processes to structured, queryable systems."]),
        ("Risk Factors", [
            "We operate in a competitive market and depend on continued innovation. Macroeconomic conditions can affect customer budgets and sales cycles. We manage concentration, security, and talent risks through deliberate controls and diversification.",
            "We maintain a conservative liquidity position to withstand volatility."]),
        ("Management Discussion and Analysis", [
            "Management evaluates the business on revenue growth, gross margin, operating margin, and cash generation. The year's results reflect deliberate trade-offs favoring durable growth and profitability.",
            "We will continue to invest in product, security, and go-to-market while protecting margins."]),
        ("Outlook", [
            "We enter the year with momentum and a strong pipeline. We expect continued revenue growth, further margin improvement, and disciplined investment in strategic initiatives.",
            "Our long-term priorities are unchanged: build durable products, earn customer trust, and compound value responsibly."]),
    ],
    "Data Privacy and Security Handbook": [
        ("Purpose and Scope", [
            "This handbook defines how the organization collects, processes, stores, and protects data. It applies to all employees, contractors, and systems that handle company or customer information.",
            "It is reviewed annually and updated as regulations and systems evolve."]),
        ("Data Classification", [
            "Data is classified as public, internal, confidential, or restricted. Handling requirements increase with sensitivity, and restricted data receives the strongest controls including encryption and strict access limits.",
            "Every dataset must have an owner responsible for its classification and lifecycle."]),
        ("Access Control", [
            "Access follows least privilege: individuals receive only the permissions their role requires. Access is reviewed periodically and revoked promptly on role change or departure.",
            "Administrative actions are logged and multi-factor authentication is required for privileged systems."]),
        ("Encryption and Storage", [
            "Sensitive data is encrypted in transit and at rest. Encryption keys are managed centrally with rotation policies, and backups are encrypted and tested for restorability.",
            "Retention schedules ensure data is not kept longer than necessary."]),
        ("Incident Response", [
            "Suspected incidents must be reported immediately. The response team triages, contains, eradicates, and recovers, then conducts a blameless postmortem to prevent recurrence.",
            "Regulatory notification obligations are assessed for every qualifying incident."]),
        ("Vendor Management", [
            "Third parties that handle company data undergo security review before onboarding and periodically thereafter. Contracts include data protection terms and breach notification requirements.",
            "Vendor access is scoped narrowly and monitored."]),
        ("Employee Responsibilities", [
            "Employees complete security training, use approved tools, protect credentials, and report concerns. Personal devices that access company data must meet baseline security requirements.",
            "Violations may result in disciplinary action."]),
        ("Compliance and Audit", [
            "The program aligns with recognized frameworks and is subject to internal and external audit. Findings are tracked to closure with clear ownership and deadlines.",
            "Evidence of controls is retained to demonstrate ongoing compliance."]),
    ],
    "Research Paper - Hybrid Retrieval at Scale": [
        ("Abstract", [
            "We study hybrid retrieval that combines lexical and dense methods for document search without a generation stage. We show that reciprocal rank fusion of BM25 and cosine similarity over compact embeddings yields robust results across heterogeneous document types while remaining inexpensive to run on commodity CPUs."]),
        ("Introduction", [
            "Retrieval systems must serve two kinds of questions: precise lookups that map to fields, and open questions best answered by meaning. Prior work often optimizes one at the expense of the other. We describe a unified approach and quantify the trade-offs.",
            "Our contributions are a practical fusion recipe, an analysis of common correctness pitfalls, and an evaluation methodology suitable for small corpora."]),
        ("Method", [
            "Documents are parsed, chunked under the embedding model's token budget, and indexed into a full-text index and a vector store. Queries run both legs and fuse the rankings. A lightweight router detects field-intent queries and answers them directly from extracted structured data.",
            "We deliberately avoid tuning fusion weights, relying instead on rank-based fusion that is stable across datasets."]),
        ("Experimental Setup", [
            "We construct a labelled set of representative queries spanning factual lookups and semantic questions. We report recall at ten and mean reciprocal rank per mode. Embeddings are computed with a small English model to keep inference on CPU.",
            "We ablate query instructions, chunk sizes, and stopword-aware snippeting."]),
        ("Results", [
            "Hybrid retrieval matches or exceeds either method alone on aggregate metrics. Semantic retrieval dominates on paraphrased questions, while lexical retrieval remains essential for identifiers and exact terms. Fusion captures the strengths of both.",
            "Snippet anchoring on content terms materially improves perceived relevance without changing ranking."]),
        ("Discussion", [
            "The dominant risks are silent correctness bugs: inverted similarity conversions, sign errors in lexical scores, and truncation from token-budget mismatches. Contract tests and a standing evaluation set catch these before they reach users.",
            "The approach generalizes to new formats because parsing and extraction are modular."]),
        ("Conclusion", [
            "A disciplined hybrid retriever provides strong, inexpensive search over messy documents. Structured extraction adds deterministic querying, and an optional generation stage can be layered on without disturbing the core.",
            "Future work includes cross-encoder reranking on CPU and richer structured extraction."]),
    ],
}

# Big decks reuse the same authored sections as slides.
DECKS = {
    "Acme FY2026 Business Review": "Acme FY2026 Annual Report",
    "Offline Document Intelligence - Tech Deck": "Whitepaper - Offline Document Intelligence",
    "Security Program Overview": "Data Privacy and Security Handbook",
}


def gen_pdfs():
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("h1b", parent=styles["Heading1"], spaceAfter=10)
    body = ParagraphStyle("bodyb", parent=styles["BodyText"], fontSize=11, leading=16, spaceAfter=10)
    for title, sections in DOCS.items():
        fname = OUT / (title.replace(" ", "_").replace("-", "") + ".pdf")
        doc = SimpleDocTemplate(str(fname), pagesize=letter,
                                topMargin=0.9 * inch, bottomMargin=0.9 * inch)
        flow = [Paragraph(title, styles["Title"]), Spacer(1, 18)]
        for heading, paras in sections:
            flow.append(Paragraph(heading, h1))
            for p in paras:
                flow.append(Paragraph(p, body))
            flow.append(Spacer(1, 8))
        doc.build(flow)


def gen_decks():
    from pptx import Presentation
    from pptx.util import Pt
    for deck_title, src in DECKS.items():
        sections = DOCS[src]
        prs = Presentation()
        # title slide
        s0 = prs.slides.add_slide(prs.slide_layouts[0])
        s0.shapes.title.text = deck_title
        s0.placeholders[1].text = "Prepared by the EasyNotes team · FY2026"
        for heading, paras in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading
            tf = slide.placeholders[1].text_frame
            tf.text = paras[0]
            for p in paras[1:]:
                para = tf.add_paragraph(); para.text = p; para.font.size = Pt(14)
            slide.notes_slide.notes_text_frame.text = (
                f"Speaker notes: expand on '{heading}'. " + paras[0][:120])
        prs.save(OUT / (deck_title.replace(" ", "_").replace("-", "") + ".pptx"))


gen_pdfs()
gen_decks()
for p in sorted(OUT.iterdir()):
    kb = p.stat().st_size // 1024
    print(f"  {p.name}  ({kb} KB)")
