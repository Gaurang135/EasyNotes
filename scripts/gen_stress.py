"""Dev tool: generate ~70 varied files (many edge cases) into a dir for stress testing.
Usage: python scripts/gen_stress.py /tmp/stress
"""
from __future__ import annotations
import sys, os, zipfile, random
from pathlib import Path

OUT = Path(sys.argv[1] if len(sys.argv) > 1 else "/tmp/stress")
OUT.mkdir(parents=True, exist_ok=True)
random.seed(7)
VENDORS = ["Acme Corp", "Globex Ltd", "Initech", "Umbrella Inc", "Soylent Co", "Hooli"]


def w(name, data):
    p = OUT / name
    p.write_bytes(data if isinstance(data, bytes) else data.encode("utf-8"))


# ---- TXT variants (prose, invoices, unicode, edge) --------------------------
w("prose_short.txt", "The coat pops open. I am free at last in the London sun.")
w("prose_long.txt", ("Chapter one. " + "The city breathes around me. " * 400))
w("unicode.txt", "Café résumé naïve — 日本語のテキスト — Ω≈ç√∫ — emoji 🎉🔥✅")
w("whitespace.txt", "     \n\n\t   \n")                       # -> should fail 'empty'
w("empty.txt", "")                                            # -> should fail 'empty'
w("rtl.txt", "مرحبا بكم في المتجر. الإجمالي 500 ريال.")
for i, v in enumerate(VENDORS):
    amt = random.randint(100, 9000)
    w(f"invoice_{i}.txt",
      f"INVOICE #{1000+i}\nVendor: {v}\nBill to: Client {i}\nInvoice date: 2026-0{i%9+1}-1{i%9}\n"
      f"Due date: 2026-0{i%9+2}-2{i%9}\nSubtotal Rs.{amt}.00\nTax Rs.{amt//10}.00\n"
      f"Total Rs.{amt+amt//10}.00\nContact billing{i}@{v.split()[0].lower()}.com or +91 98{i}6543210\n"
      f"https://{v.split()[0].lower()}.com/pay")
w("keyvalues.txt", "Name: Alice\nRole: Engineer\nStart: 2025-01-02\nSalary: Rs.90000\nEmail: alice@corp.io")
w("numbers_only.txt", "1234 5678 91011 " * 50)
w("code.txt", "def f(x):\n    return x*2  # doubles\nprint(f(21))  # 42")
w("mixed_fields.txt", "Reach me at bob@x.io, +1 415 555 0100, https://x.io on 12/31/2026 for $2,500.")

# ---- Markdown ---------------------------------------------------------------
w("md_headings.md", "# Title\n## Section A\nSome text.\n## Section B\nMore text about invoices.")
w("md_nested.md", "# Doc\n### Deep\ncontent\n#### Deeper\nmore")
w("md_only_heading.md", "# Just a heading")
w("md_table.md", "# Report\n\n| a | b |\n|---|---|\n| 1 | 2 |\n")

# ---- CSV variants (delimiters, ragged, quoted, edge) ------------------------
w("csv_normal.csv", "region,amount,date\nnorth,1200,2026-01-05\nsouth,300,2026-02-11\n")
w("csv_ragged.csv", "a,b,c\n1,2\n3,4,5,6\n,,\n7,8,9\n")
w("csv_quoted.csv", 'name,note\n"Smith, John","said ""hi"""\n"Doe, Jane","two, commas"\n')
w("csv_semicolon.csv", "x;y;z\n1;2;3\n4;5;6\n")
w("csv_tab.csv", "p\tq\tr\n1\t2\t3\n")
w("csv_one_col.csv", "id\n1\n2\n3\n")
w("csv_no_header_numeric.csv", "1,2,3\n4,5,6\n7,8,9\n")
w("csv_unicode.csv", "naïve,montante,fëcha\ncafé,1200,2026-01-05\n")
w("csv_empty.csv", "")                                        # -> fail
w("csv_headers_only.csv", "a,b,c\n")                          # header, no rows
w("csv_big.csv", "id,val\n" + "\n".join(f"{i},{i*3}" for i in range(2000)))
w("csv_mixed_types.csv", "id,amount,flag,date\n1,10.5,true,2026-01-01\n2,x,false,notadate\n")
w("csv_amounts.csv", "item,price\nwidget,Rs.1200\ngadget,Rs.3400\ngizmo,Rs.560\n")

# ---- office + pdf via libs -------------------------------------------------
def gen_office():
    from docx import Document
    from pptx import Presentation
    from openpyxl import Workbook
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from pypdf import PdfReader, PdfWriter

    # DOCX
    for i in range(4):
        d = Document()
        d.add_heading(f"Report {i}", level=1)
        d.add_paragraph(f"Vendor: {VENDORS[i]}. Total Rs.{(i+1)*1000}. contact{i}@corp.com")
        t = d.add_table(rows=3, cols=3)
        for r in range(3):
            for c in range(3):
                t.cell(r, c).text = f"r{r}c{c}"
        d.save(OUT / f"docx_{i}.docx")
    Document().save(OUT / "docx_empty.docx")                  # empty -> fail

    # PPTX
    for i in range(4):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"Deck {i}: {VENDORS[i]}"
        s.placeholders[1].text = f"Total Rs.{(i+1)*500}. deck{i}@corp.com on 2026-0{i+1}-15"
        s.notes_slide.notes_text_frame.text = f"notes for deck {i}"
        prs.save(OUT / f"pptx_{i}.pptx")
    Presentation().save(OUT / "pptx_empty.pptx")              # no text -> fail

    # XLSX
    for i in range(4):
        wb = Workbook(); ws = wb.active; ws.title = f"S{i}"
        ws.append(["id", "amount", "total"])
        for r in range(1, 6):
            ws.append([r, r * (i + 1) * 10, f"=B{r+1}*2"])
        if i == 0:
            ws2 = wb.create_sheet("Second"); ws2.append(["k", "v"]); ws2.append(["a", 1])
        wb.save(OUT / f"xlsx_{i}.xlsx")
    Workbook().save(OUT / "xlsx_empty.xlsx")                  # empty -> fail

    # PDF
    for i in range(3):
        base = OUT / f"pdf_{i}.pdf"
        c = canvas.Canvas(str(base), pagesize=letter)
        for pg in range(i + 1):
            c.drawString(72, 720, f"PDF {i} page {pg+1}: invoice total Rs.{(i+1)*111}. pay@corp.com")
            c.showPage()
        c.save()
    # owner-locked (must parse) + user-locked (must fail encrypted)
    src = OUT / "pdf_0.pdf"
    r = PdfReader(str(src)); wtr = PdfWriter()
    for p in r.pages:
        wtr.add_page(p)
    wtr.encrypt(user_password="", owner_password="o")
    (OUT / "pdf_owner_locked.pdf").write_bytes(b"")
    with open(OUT / "pdf_owner_locked.pdf", "wb") as f:
        wtr.write(f)
    r2 = PdfReader(str(src)); w2 = PdfWriter()
    for p in r2.pages:
        w2.add_page(p)
    w2.encrypt(user_password="secret", owner_password="o")
    with open(OUT / "pdf_user_locked.pdf", "wb") as f:
        w2.write(f)


gen_office()

# ---- corrupt / malicious / wrong-type edge cases ---------------------------
w("corrupt.pdf", b"%PDF-1.4 not really a pdf \x00\x01\x02")   # corrupt -> fail
w("corrupt.docx", b"PK\x03\x04 broken zip")                   # corrupt -> fail
w("fake.docx", b"just text pretending")                        # not a zip -> sniff rejects
z = OUT / "renamed_zip.docx"                                    # real zip, not docx
with zipfile.ZipFile(z, "w") as zf:
    zf.writestr("junk.txt", "hello")
w("junk.exe", b"MZ\x90\x00\x03binary junk")                    # wrong type -> 415
w("noext", b"some bytes without extension")                    # unknown -> 415
# zip bomb-ish xlsx (high ratio sharedStrings)
bomb = OUT / "bomb.xlsx"
with zipfile.ZipFile(bomb, "w", zipfile.ZIP_DEFLATED) as zf:
    zf.writestr("xl/workbook.xml", "<x/>")
    zf.writestr("xl/sharedStrings.xml", b"0" * (60 * 1024 * 1024))  # -> UnsafeArchive

files = sorted(os.listdir(OUT))
print(f"generated {len(files)} files in {OUT}")
