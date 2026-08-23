"""Dev-only: generate the committed binary test fixtures. Run once via the venv."""
from pathlib import Path

FIX = Path("tests/fixtures")
FIX.mkdir(parents=True, exist_ok=True)


def gen_docx():
    from docx import Document
    d = Document()
    d.add_heading("Quarterly Report", level=1)
    d.add_paragraph("Revenue grew across all payment methods this quarter.")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Method"; t.cell(0, 1).text = "Amount"
    t.cell(1, 0).text = "UPI"; t.cell(1, 1).text = "1200"
    d.save(FIX / "sample.docx")


def gen_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Payments Overview"
    slide.placeholders[1].text = "Refunds and chargebacks summary"
    slide.notes_slide.notes_text_frame.text = "Speaker notes: mention the refund SLA"
    prs.save(FIX / "sample.pptx")


def gen_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["id", "amount", "total"])
    ws.append([1, 100, "=B2*2"])
    ws.append([2, 250, "=B3*2"])
    wb.save(FIX / "sample.xlsx")


def gen_pdfs():
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    base = FIX / "simple.pdf"
    c = canvas.Canvas(str(base), pagesize=letter)
    c.drawString(72, 720, "This is a simple PDF about invoices and payments.")
    c.showPage()
    c.save()

    from pypdf import PdfReader, PdfWriter
    # owner-password-only (empty user password) -> must still parse
    r = PdfReader(str(base)); w = PdfWriter()
    for p in r.pages:
        w.add_page(p)
    w.encrypt(user_password="", owner_password="owner")
    with open(FIX / "owner_locked.pdf", "wb") as f:
        w.write(f)
    # real user password -> must raise EncryptedFileError
    r2 = PdfReader(str(base)); w2 = PdfWriter()
    for p in r2.pages:
        w2.add_page(p)
    w2.encrypt(user_password="secret", owner_password="owner")
    with open(FIX / "user_locked.pdf", "wb") as f:
        w2.write(f)


if __name__ == "__main__":
    gen_docx(); gen_pptx(); gen_xlsx(); gen_pdfs()
    print("fixtures generated:", sorted(p.name for p in FIX.iterdir()))
